from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from typing import Dict, List, Optional, Tuple
import os

class PPTXParser:
    """Parse and analyze existing PowerPoint files"""
    
    def __init__(self, pptx_path: str):
        """Initialize parser with a PPTX file"""
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)
    
    def extract_colors(self) -> List[str]:
        """Extract color palette from presentation"""
        colors = set()
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                try:
                    
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.color.type == 1:  # RGB color
                                    rgb = run.font.color.rgb
                                    colors.add(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                    
                    
                    if hasattr(shape, 'fill') and shape.fill.type == 1:  # Solid fill
                        rgb = shape.fill.fore_color.rgb
                        colors.add(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                except:
                    pass
        
        return list(colors)[:10]  
    
    def extract_fonts(self) -> List[str]:
        """Extract fonts used in presentation"""
        fonts = set()
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                fonts.add(run.font.name)
        
        return list(fonts)
    
    def extract_layouts(self) -> List[str]:
        """Extract layout names from presentation"""
        layouts = []
        for layout in self.prs.slide_layouts:
            layouts.append(layout.name)
        return layouts
    
    def get_slide_structure(self) -> List[Dict]:
        """Extract structure of all slides"""
        slides = []
        
        for idx, slide in enumerate(self.prs.slides):
            slide_data = {
                'index': idx,
                'layout': slide.slide_layout.name,
                'shapes': []
            }
            
            for shape in slide.shapes:
                shape_data = {
                    'type': shape.shape_type,
                    'has_text': shape.has_text_frame
                }
                
                if shape.has_text_frame:
                    text = '\n'.join([paragraph.text for paragraph in shape.text_frame.paragraphs])
                    shape_data['text'] = text
                
                slide_data['shapes'].append(shape_data)
            
            slides.append(slide_data)
        
        return slides
    
    def get_style_profile(self) -> Dict:
        """Get comprehensive style profile of presentation"""
        return {
            'colors': self.extract_colors(),
            'fonts': self.extract_fonts(),
            'layouts': self.extract_layouts(),
            'slide_count': len(self.prs.slides),
            'slide_width': self.prs.slide_width,
            'slide_height': self.prs.slide_height
        }
    
    def extract_template_metadata(self) -> Dict:
        """Extract metadata for template reuse"""
        return {
            'style_profile': self.get_style_profile(),
            'slide_structure': self.get_slide_structure(),
            'master_slides': len(self.prs.slide_master.slide_layouts)
        }
