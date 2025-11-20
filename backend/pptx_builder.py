from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import Dict, List, Optional
import os
from backend.config import config

class PPTXBuilder:
    """Build PowerPoint presentations programmatically"""
    
    def __init__(self, template_path: Optional[str] = None):
        """Initialize builder with optional template"""
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        self.current_style = {
            'primary_color': RGBColor(31, 71, 136),  # Default blue
            'secondary_color': RGBColor(74, 144, 226),
            'accent_color': RGBColor(255, 255, 255),
            'font_name': config.DEFAULT_FONT
        }
    
    def set_style(self, style: Dict):
        """Set presentation style from dictionary"""
        if 'primary_color' in style:
            self.current_style['primary_color'] = self._hex_to_rgb(style['primary_color'])
        if 'secondary_color' in style:
            self.current_style['secondary_color'] = self._hex_to_rgb(style['secondary_color'])
        if 'accent_color' in style:
            self.current_style['accent_color'] = self._hex_to_rgb(style['accent_color'])
        if 'font_suggestion' in style:
            self.current_style['font_name'] = style['font_suggestion']
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    
    def _get_layout_by_name(self, layout_name: str):
        """Get slide layout by name"""
        for layout in self.prs.slide_layouts:
            if layout_name.lower() in layout.name.lower():
                return layout
        
        return self.prs.slide_layouts[0]
    
    def add_title_slide(self, title: str, subtitle: str = ""):
        """Add a title slide"""
        layout = self._get_layout_by_name('Title')
        slide = self.prs.slides.add_slide(layout)
        
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._apply_text_style(slide.shapes.title, config.TITLE_FONT_SIZE, bold=True)
        
        
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
            self._apply_text_style(slide.placeholders[1], config.DEFAULT_FONT_SIZE + 4)
        
        return slide
    
    def add_content_slide(self, title: str, content: List[str]):
        """Add a slide with title and bullet points"""
        layout = self._get_layout_by_name('Content')
        slide = self.prs.slides.add_slide(layout)
        
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._apply_text_style(slide.shapes.title, config.HEADING_FONT_SIZE, bold=True)
        
       
        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            
            for item in content:
                p = text_frame.add_paragraph()
                p.text = item
                p.level = 0
                self._apply_paragraph_style(p, config.DEFAULT_FONT_SIZE)
        
        return slide
    
    def add_section_header(self, title: str):
        """Add a section header slide"""
        layout = self._get_layout_by_name('Section')
        if not layout:
            layout = self._get_layout_by_name('Title Only')
        
        slide = self.prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._apply_text_style(slide.shapes.title, config.TITLE_FONT_SIZE, bold=True)
        
        return slide
    
    def add_two_column_slide(self, title: str, left_content: List[str], right_content: List[str]):
        """Add a slide with two columns of content"""
        layout = self._get_layout_by_name('Two Content')
        if not layout:
            layout = self._get_layout_by_name('Blank')
        
        slide = self.prs.slides.add_slide(layout)
        
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._apply_text_style(slide.shapes.title, config.HEADING_FONT_SIZE, bold=True)
        
    
        if len(slide.placeholders) > 2:
            
            left_frame = slide.placeholders[1].text_frame
            left_frame.clear()
            for item in left_content:
                p = left_frame.add_paragraph()
                p.text = item
                self._apply_paragraph_style(p, config.DEFAULT_FONT_SIZE)
            
           
            right_frame = slide.placeholders[2].text_frame
            right_frame.clear()
            for item in right_content:
                p = right_frame.add_paragraph()
                p.text = item
                self._apply_paragraph_style(p, config.DEFAULT_FONT_SIZE)
        else:
            
            left_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2), Inches(4.5), Inches(4)
            )
            left_frame = left_box.text_frame
            for item in left_content:
                p = left_frame.add_paragraph()
                p.text = item
                self._apply_paragraph_style(p, config.DEFAULT_FONT_SIZE)
            
            right_box = slide.shapes.add_textbox(
                Inches(5.5), Inches(2), Inches(4.5), Inches(4)
            )
            right_frame = right_box.text_frame
            for item in right_content:
                p = right_frame.add_paragraph()
                p.text = item
                self._apply_paragraph_style(p, config.DEFAULT_FONT_SIZE)
        
        return slide
    
    def add_image_slide(self, title: str, image_path: str, caption: str = ""):
        """Add a slide with an image"""
        layout = self._get_layout_by_name('Picture')
        if not layout:
            layout = self._get_layout_by_name('Blank')
        
        slide = self.prs.slides.add_slide(layout)
        
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._apply_text_style(slide.shapes.title, config.HEADING_FONT_SIZE, bold=True)
        
       
        if os.path.exists(image_path):
            left = Inches(2)
            top = Inches(2)
            height = Inches(4)
            slide.shapes.add_picture(image_path, left, top, height=height)
        
       
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(1), Inches(6.5), Inches(8), Inches(0.5)
            )
            caption_frame = caption_box.text_frame
            p = caption_frame.paragraphs[0]
            p.text = caption
            p.alignment = PP_ALIGN.CENTER
            self._apply_paragraph_style(p, config.DEFAULT_FONT_SIZE - 2)
        
        return slide
    
    def _apply_text_style(self, shape, font_size: int, bold: bool = False):
        """Apply text styling to a shape"""
        if not shape.has_text_frame:
            return
        
        for paragraph in shape.text_frame.paragraphs:
            self._apply_paragraph_style(paragraph, font_size, bold)
    
    def _apply_paragraph_style(self, paragraph, font_size: int, bold: bool = False):
        """Apply styling to a paragraph"""
        for run in paragraph.runs:
            run.font.name = self.current_style['font_name']
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = self.current_style['primary_color']
    
    def build_from_structure(self, structure: Dict):
        """Build presentation from LLM-generated structure"""
        slides = structure.get('slides', [])
        
        for slide_data in slides:
            slide_type = slide_data.get('type', 'Title and Content')
            
            if slide_type == 'Title Slide':
                self.add_title_slide(
                    slide_data.get('title', ''),
                    slide_data.get('subtitle', '')
                )
            elif slide_type == 'Section Header':
                self.add_section_header(slide_data.get('title', ''))
            elif slide_type == 'Title and Content':
                self.add_content_slide(
                    slide_data.get('title', ''),
                    slide_data.get('content', [])
                )
            elif slide_type == 'Two Content':
                self.add_two_column_slide(
                    slide_data.get('title', ''),
                    slide_data.get('left_content', []),
                    slide_data.get('right_content', [])
                )
    
    def save(self, output_path: str):
        """Save presentation to file"""
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self.prs.save(output_path)
        return output_path
